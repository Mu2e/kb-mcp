"""PowerPoint parser for extracting text from PPTX files."""

import logging
from pathlib import Path

from .parser_base import BaseParser
from .text_utils import slides_format_as_markdown

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """Parser for PowerPoint presentations (PPTX)."""

    def extract_text(self) -> str:
        """Extract raw text from PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
            return ""

        try:
            prs = Presentation(self.file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {self.file_path}: {e}")
            return ""

    def get_text(self) -> str:
        """Get formatted text from PowerPoint presentation.
        
        Applies slide formatting as markdown, then standard cleaning.
        """
        text = self.extract_text()
        # Apply slide-specific formatting
        text = slides_format_as_markdown(text)
        # Then apply standard cleaning
        from .text_utils import clean_text
        return clean_text(text)


